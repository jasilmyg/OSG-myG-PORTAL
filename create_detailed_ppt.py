from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Function to add a title slide
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

# Function to add a bullet point slide
def add_bullet_slide(prs, title, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = title
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        if isinstance(point, list):
            for subpoint in point:
                p = tf.add_paragraph()
                p.text = subpoint
                p.level = 1

# Slide 1: Title
add_title_slide(
    prs, 
    "Q2 2026 Detailed Accomplishments", 
    "Comprehensive Review of Development Work\nApril, May & June 2026"
)

# Slide 2: Executive Summary
add_bullet_slide(
    prs,
    "Executive Summary",
    [
        "Architected the backend for OSG-myG-PORTAL (WhatsApp APIs, OSID tracking, Claims processing).",
        "Developed and deployed an Enterprise AI Agent powered by LLMs (NVIDIA/Llama) and Machine Learning.",
        "Engineered the Enterprise Retail Analytics Dashboard & automated 'SHE START' live-syncing system.",
        "Executed heavy Database Administration (refreshing 26+ materialized views, optimizing 12.6M+ row imports).",
        "Delivered technical documentation and bespoke UI projects (Glassmorphism web apps)."
    ]
)

# Slide 3: OSG-myG-PORTAL (Backend & WhatsApp)
add_bullet_slide(
    prs,
    "OSG-myG-PORTAL Architecture",
    [
        "Engineered a robust claims processing workflow with real-time WhatsApp API integration.",
        "Developed custom python modules for OSID data parsing (import_osid) and forensic data investigation (trace_9895.py).",
        "Implemented Google Apps Script integrations for synced data processing.",
        "Built fault-tolerant data pipelines with audit status tracking and schema validation.",
        "Restored broken workflows and optimized Postman API fetching via Python."
    ]
)

# Slide 4: Enterprise AI Agent & ML Pipelines
add_bullet_slide(
    prs,
    "Enterprise AI Agent & ML Pipelines",
    [
        "Architected a multi-layered AI Agent including a SQL Agent and AI Analyst for the Loyalty Portal.",
        "Integrated high-performance LLMs via NVIDIA API (Llama models) to answer complex business questions in natural language.",
        "Configured dynamic prompt engineering for deep context awareness and fast execution.",
        "Deployed ML models (Random Forest, MLPRegressor) and PyTorch concepts for dormant customer forecasting.",
        "Built the 'MalayalamCalendarFeaturizer' to factor local Kerala holidays (e.g., Onam) into predictive scoring."
    ]
)

# Slide 5: Enterprise Retail Analytics & Data Management
add_bullet_slide(
    prs,
    "Analytics & Database Management",
    [
        "Refactored materialized views and ETL pipelines in PostgreSQL (handling 12.6+ million rows) to eliminate cross-joins.",
        "Automated the DB Manager upload flow with custom Python scripts to bypass timeouts and sanitize anomalies (e.g., SMC/EI).",
        "Optimized LocMemCache and Redis layers in Django to deliver sub-second data loading.",
        "Replaced standard Excel generators with high-speed Rust-based calamine engines.",
        "Built high-speed DataTables grids for month-wise, quarterly, and FY comparisons."
    ]
)

# Slide 6: SHE START Applicant Tracking Dashboard
add_bullet_slide(
    prs,
    "SHE START Applicant Tracking System",
    [
        "Created a live Google Sheets sync engine (gspread) with a 25-second automated refresh.",
        "Integrated interactive, inline cell-editing with silent background saving to PostgreSQL.",
        "Programmed a complex Final Decision Matrix, dynamically dropping extreme scores across 6 panelists to compute weighted averages.",
        "Automated decision badges (Strong Final Selection, Recommended for Top 10, Waitlist).",
        "Configured secure, role-based access control to isolate views from standard users."
    ]
)

# Slide 7: Technical Documentation & Frontend Projects
add_bullet_slide(
    prs,
    "Documentation & Frontend Engineering",
    [
        "Performed deep source-code analysis to generate comprehensive GitHub README.md documentation.",
        "Engineered an Ultra-Premium Birthday Website featuring a black & gold luxury theme, glassmorphism UI, and cinematic animations.",
        "Upgraded Plotly Dashboards with dynamic UI elements like Probability Gauges, Semi-Donuts, and Dormancy Risk Meters.",
        "Managed environment variables, SSH tunneling to DigitalOcean, and GitHub version control."
    ]
)

# Save the presentation
output_path = "c:\\Users\\jasil_myg\\Desktop\\OSG-myG-PORTAL-mainnnnn - Copy\\OSG_Portal_WorkSummary_Q2_2026_Detailed.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
