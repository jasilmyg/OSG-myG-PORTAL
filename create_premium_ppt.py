from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Slide layouts
TITLE_SLIDE = 0
BULLET_SLIDE = 1
SECTION_HEADER = 2
TITLE_ONLY = 5

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

def add_section_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    tf = body.text_frame
    tf.clear()
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph()
        if isinstance(point, list):
            # sub-bullets
            for subpoint in point:
                sp = tf.add_paragraph()
                sp.text = subpoint
                sp.level = 1
                sp.font.size = Pt(18)
        else:
            p.text = point
            p.level = 0
            p.font.size = Pt(22)
            # Add some spacing before major points
            if i > 0:
                p.space_before = Pt(12)

# --- SLIDES ---

# 1. Title
add_title_slide(prs, "Q2 2026 Engineering Review", "Development, AI, and Database Architecture\nApril - June 2026")

# 2. Executive Summary
add_content_slide(prs, "Executive Summary", [
    "Delivered 5 major technical projects across multiple environments in Q2 2026.",
    "Integrated Enterprise AI into analytics utilizing NVIDIA LLMs.",
    "Engineered robust WhatsApp API pipelines for Onsitego claims tracking.",
    "Optimized heavy PostgreSQL architectures handling over 12.6 million records.",
    "Developed highly scalable OTT and Applicant Tracking web applications."
])

# 3. Section: OSG-myG-PORTAL
add_section_slide(prs, "Project 1: OSG-myG-PORTAL")
add_content_slide(prs, "Claims & WhatsApp Integration", [
    "WhatsApp API Routing: Engineered automated notification pipelines via Telinfy / GreenAds Global APIs.",
    "Data Forensics: Built custom modules (trace_9895.py, investigate_claim.py) to resolve workflow anomalies.",
    "ETL Pipelines: Automated the ingestion of legacy OSID datasets while standardizing date/time conflicts.",
    "Google Apps Script: Bridged live Google Sheets with the Django backend for automated audits."
])

# 4. Section: AI & Machine Learning
add_section_slide(prs, "Project 2: Enterprise AI Agent & ML")
add_content_slide(prs, "Conversational BI & Forecasting", [
    "Enterprise AI Analyst: Architected a SQL-translating Agent powered by NVIDIA/Llama LLMs.",
    "Predictive Forecasting: Deployed Scikit-Learn (Random Forest, MLPRegressor) for dormant customer reactivation metrics.",
    "Custom Featurization: Built the 'MalayalamCalendarFeaturizer' to factor local Kerala holidays (e.g., Onam) into models.",
    "Optimized API Latency: Refined prompt engineering to prevent database schema hallucinations and reduce execution time."
])

# 5. Section: Database & Analytics
add_section_slide(prs, "Project 3: Loyalty Portal Architecture")
add_content_slide(prs, "Database & Performance Optimization", [
    "PostgreSQL Optimization: Repaired 'mv_yearly_cohort' cross-joins and automated concurrent refreshing for 26 massive materialized views.",
    "High-Speed Caching: Integrated Redis & LocMemCache, reducing 12.6M+ row query loads from minutes to milliseconds.",
    "DB Manager Security: Automated the scrubbing of anomalous internal store data (SMC/EI) directly during batch uploads.",
    "Excel Export Engine: Deployed the Rust-based 'calamine' engine to accelerate report generation by 5-10x."
])

# 6. Section: SHE START
add_section_slide(prs, "Project 4: 'SHE START' Dashboard")
add_content_slide(prs, "Applicant Evaluation System", [
    "Live-Sync Architecture: Engineered a 25-second automated 'gspread' sync mirroring Google Sheets to the Django portal.",
    "Complex Math Logic: Built an anti-bias scoring algorithm (dropping highest & lowest panelist scores before averaging).",
    "Interactive UI: Implemented inline cell-editing with silent background saving to Postgres.",
    "Role-Based Access: Secured the portal to isolate sensitive evaluation metrics from standard viewers."
])

# 7. Section: Bigg Boss / FoneFlix
add_section_slide(prs, "Project 5: OTT Registration Portals")
add_content_slide(prs, "Bigg Boss & FoneFlix Web Apps", [
    "OAuth 2.0 Integration: Built a custom Google Drive Auth flow, bypassing server storage limits for large audition video uploads.",
    "Premium OTT Design: Engineered a 40/60 split UI featuring dark glassmorphism, neon glows, and 3D floating particles.",
    "Dynamic Rebranding: Successfully pivoted the entire codebase from Bigg Boss (Navy/Pink) to FoneFlix (Dark Browns/Orange).",
    "Secure Backend: Flask-based architecture handling user consent, file processing, and robust data storage."
])

# 8. Business Impact
add_content_slide(prs, "Value Delivered in Q2", [
    "Performance: Reduced multi-million row DB query times dramatically via Redis and materialized view optimization.",
    "Automation: Replaced manual Excel cross-referencing with AI-powered conversational querying.",
    "Data Integrity: Hardened claims workflows, ensuring no dropped WhatsApp notifications or anomalous sales records.",
    "Scalability: OAuth 2.0 implementations removed infrastructure bottlenecks for large media uploads."
])

# Save
output_path = r"c:\Users\jasil_myg\Desktop\OSG-myG-PORTAL-mainnnnn - Copy\Q2_2026_Premium_Work_Report.pptx"
prs.save(output_path)
print(f"Premium Presentation saved to {output_path}")
