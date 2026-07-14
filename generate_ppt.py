"""
OSG myG Portal - Q2 2026 Work Summary
Most Accurate Version - Based on Git History + Full Chat Transcript Analysis
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# Premium Palette
BG    = RGBColor(0x06, 0x09, 0x14)
CARD  = RGBColor(0x0C, 0x15, 0x28)
CARD2 = RGBColor(0x11, 0x1C, 0x35)
STRIP = RGBColor(0x08, 0x11, 0x20)
BLUE  = RGBColor(0x38, 0x8B, 0xFF)
PURP  = RGBColor(0xA7, 0x6A, 0xFF)
PINK  = RGBColor(0xF0, 0x4A, 0x9D)
GREE  = RGBColor(0x0D, 0xC8, 0x8A)
AMBE  = RGBColor(0xFF, 0xA5, 0x0A)
CYAN  = RGBColor(0x00, 0xD4, 0xFF)
RED   = RGBColor(0xFF, 0x5C, 0x5C)
WHITE = RGBColor(0xEE, 0xF2, 0xFF)
MUTED = RGBColor(0x7A, 0x8F, 0xC2)
BORD  = RGBColor(0x1A, 0x30, 0x55)

def fix_rgb(c):
    # RGBColor components
    r = (c >> 16) & 0xFF
    g = (c >> 8) & 0xFF
    b = c & 0xFF
    return r, g, b

MUTED = RGBColor(0x7A, 0x8F, 0xC2)

def Rc(sl, l, t, w, h, fill=None, line=None, lw=0.5):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def Tx(sl, text, l, t, w, h, sz=12, bold=False, italic=False,
       col=WHITE, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = col; r.font.name = "Calibri"
    return tb

def bullets(sl, items, l, t, w, h, sz=11, col=MUTED, spacing=3):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(spacing)
        r = p.add_run(); r.text = item
        r.font.size = Pt(sz); r.font.color.rgb = col; r.font.name = "Calibri"

def page_bg(sl): Rc(sl, 0, 0, 13.33, 7.5, fill=BG)

def top_bar(sl, clr):
    Rc(sl, 0, 0, 13.33, 0.055, fill=clr)

def hdr(sl, title, sub, clr):
    Rc(sl, 0, 0, 13.33, 1.55, fill=CARD)
    top_bar(sl, clr)
    Tx(sl, title, 0.45, 0.1, 11, 0.82, sz=30, bold=True, col=WHITE)
    Tx(sl, sub, 0.45, 0.9, 12.5, 0.45, sz=12, col=MUTED)

def feat_card(sl, title, items, l, t, w, h, clr=BLUE, tsz=12, isz=10.5):
    Rc(sl, l, t, w, h, fill=CARD2, line=BORD, lw=0.3)
    Rc(sl, l, t, 0.055, h, fill=clr)
    Rc(sl, l, t, w, 0.38, fill=STRIP)
    Tx(sl, title, l+0.15, t+0.06, w-0.2, 0.3, sz=tsz, bold=True, col=clr)
    bullets(sl, items, l+0.15, t+0.44, w-0.22, h-0.5, sz=isz, col=MUTED, spacing=4)

def kpi(sl, val, lbl, l, t, clr=BLUE):
    Rc(sl, l, t, 2.75, 1.38, fill=CARD2, line=BORD, lw=0.3)
    Rc(sl, l, t+1.28, 2.75, 0.1, fill=clr)
    Tx(sl, val, l, t+0.06, 2.75, 0.78, sz=40, bold=True, col=clr, align=PP_ALIGN.CENTER)
    Tx(sl, lbl, l, t+0.88, 2.75, 0.38, sz=9.5, col=MUTED, align=PP_ALIGN.CENTER)

def btm(sl, pairs, clr):
    Rc(sl, 0, 6.52, 13.33, 0.98, fill=CARD)
    Rc(sl, 0, 6.52, 13.33, 0.04, fill=clr)
    w = 13.33 / len(pairs)
    for i, (v, l) in enumerate(pairs):
        x = i * w
        Tx(sl, v, x, 6.6, w, 0.5, sz=20, bold=True, col=clr, align=PP_ALIGN.CENTER)
        Tx(sl, l, x, 7.1, w, 0.3, sz=9, col=MUTED, align=PP_ALIGN.CENTER)

# =============================================================
# SLIDE 1 - TITLE
# =============================================================
sl = prs.slides.add_slide(BLANK)
page_bg(sl)

# Left panel
Rc(sl, 0, 0, 6.7, 7.5, fill=RGBColor(0x09, 0x12, 0x22))
for y, c in [(0, BLUE), (0.06, PURP), (0.12, PINK)]:
    Rc(sl, 0, y, 6.7, 0.05, fill=c)

# Vertical pattern on right
for i in range(20):
    Rc(sl, 6.85 + i * 0.33, 0, 0.018, 7.5, fill=RGBColor(0x10, 0x1E, 0x38))

# Diagonal glow circle (right)
for sz, clr_r, clr_g, clr_b in [
    (4.5, 0x00, 0x4A, 0xFF), (3.2, 0x00, 0x3A, 0xCC), (1.8, 0x00, 0x6E, 0xFF)
]:
    s = sl.shapes.add_shape(9, Inches(8.0), Inches(7.5/2 - sz/2.2),
                             Inches(sz * 1.2), Inches(sz))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(clr_r, clr_g, clr_b)
    s.line.fill.background()

# Left content
Tx(sl, "APRIL  |  MAY  |  JUNE  2026", 0.5, 0.82, 5.9, 0.5,
   sz=13, bold=True, col=BLUE)
Tx(sl, "OSG myG", 0.4, 1.38, 5.8, 1.1, sz=54, bold=True, col=WHITE)
Tx(sl, "Portal", 0.4, 2.4, 5.8, 0.95, sz=54, bold=True, col=BLUE)
Tx(sl, "Q2 2026  Work Summary", 0.4, 3.4, 5.8, 0.58, sz=22, bold=True, col=PURP)
Rc(sl, 0.4, 4.1, 3.5, 0.04, fill=BLUE)
Tx(sl, "Prepared for Manager Review", 0.4, 4.28, 5.8, 0.4, sz=13, italic=True, col=MUTED)

# Right info cards
for i, (lbl, val, clr) in enumerate([
    ("PROJECT",     "OSG myG Customer Portal",         BLUE),
    ("ROLE",        "Full Stack Python Developer",      PURP),
    ("STACK",       "Flask · PostgreSQL · JS · CSS",   GREE),
    ("DEPLOYED ON", "Render Cloud (Production Live)",   AMBE),
]):
    y = 1.05 + i * 1.52
    Rc(sl, 7.2, y, 5.8, 1.35, fill=CARD, line=BORD, lw=0.3)
    Rc(sl, 7.2, y, 0.06, 1.35, fill=clr)
    Tx(sl, lbl, 7.38, y + 0.1, 5.3, 0.3, sz=9, bold=True, col=clr)
    Tx(sl, val, 7.38, y + 0.48, 5.3, 0.62, sz=16, bold=True, col=WHITE)

for y, c in [(7.42, PINK), (7.46, PURP), (7.5, BLUE)]:
    Rc(sl, 0, y, 13.33, 0.05, fill=c)


# =============================================================
# SLIDE 2 - APRIL 2026 | PART 1: WHAT WAS BUILT
# =============================================================
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
hdr(sl, "April 2026  —  Complete Portal Built from Scratch  (1 / 2)",
    "Git Commit: 72334fd  |  Date: April 20, 2026  |  35 files added  |  12,776 lines of code written", BLUE)

feat_card(sl, "Backend: Flask Application  (app.py  =  2,817 lines)", [
    ">>  /dashboard         Admin portal: claim table, date filter, search, status filter, aging view",
    ">>  /submit-claim      Multi-product claim form — validates, saves, sends email, pushes to Sheets",
    ">>  /update-claim      Updates status, notes, follow-up, workflow steps, dates in DB + Sheets",
    ">>  /claim/<id>        JSON API — returns full claim data for modal population",
    ">>  /analytics         Returns KPI counts, TAT data, replacement progress for dashboard",
    ">>  /customer-care/<osid>   Public page — customers check claim status using OSID number",
    ">>  /login  /logout    Session-based auth with role-based access control (RBAC)",
    ">>  /forgot-password   Password reset via token sent to registered email (15 min expiry)",
    ">>  /submit-claim GET  Renders the claim submission form for admin use",
], 0.35, 1.62, 6.05, 5.5, clr=BLUE, tsz=12, isz=10)

feat_card(sl, "Frontend: 8 HTML Templates", [
    "dashboard.html          1,203 lines  - Full portal dashboard UI",
    "analytics.html            639 lines  - Analytics dashboard with KPIs",
    "claim_status.html         538 lines  - Public customer status page",
    "reports_tools.html        656 lines  - Reports & tools section",
    "submit.html               173 lines  - Claim submission form",
    "workflow_enhanced.html    323 lines  - Replacement workflow section",
    "claim_detail.html         133 lines  - Claim detail popup",
    "login.html                280 lines  - Login + forgot password",
], 6.55, 1.62, 6.43, 3.28, clr=PURP, tsz=12, isz=10)

feat_card(sl, "JavaScript: 3 Modules  (2,438 lines total)", [
    "analytics.js    1,203 lines - Charts, workflow modal, TAT calc, KPI logic",
    "script.js         751 lines - Edit modal, status update, follow-up, tabs",
    "script_v2.js      484 lines - Enhanced UX, claim detail interactions",
], 6.55, 5.05, 6.43, 2.07, clr=CYAN, tsz=12, isz=10)

btm(sl, [
    ("35", "Files Created"), ("12,776", "Lines of Code"),
    ("2,817", "app.py Lines"), ("8", "HTML Templates"), ("2,438", "JS Lines"),
], BLUE)


# =============================================================
# SLIDE 3 - APRIL 2026 | PART 2: CSS, DATA, DEPLOYMENT
# =============================================================
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
hdr(sl, "April 2026  —  Complete Portal Built from Scratch  (2 / 2)",
    "CSS Design System  |  Google Sheets Integration  |  Reference Data  |  Production Deployment", BLUE)

feat_card(sl, "CSS Design System: 5 Files  (2,725 lines total)", [
    "style.css                 1,210 lines  - Full portal design: layout, tables, cards, modals",
    "analytics.css               624 lines  - Analytics dashboard + chart styling",
    "claim_status.css            559 lines  - Customer-facing status page styling",
    "workflow_enhanced.css       321 lines  - Replacement workflow UI components",
    "grid.css                     11 lines  - Grid layout utility classes",
], 0.35, 1.62, 6.05, 3.0, clr=AMBE, tsz=12, isz=10)

feat_card(sl, "Google Sheets Integration & Reference Data", [
    "google_apps_script.js    274 lines  - Apps Script for bidirectional Sheet sync",
    "Onsitego OSID Database   15.8 MB    - Full OSID data file (updated Jan 2026)",
    "Future Store List.xlsx              - Store classification (future / existing)",
    "RBM, BDM, Branch.xlsx              - Region / Branch / Manager mapping data",
    "myG All Store.xlsx                  - Complete store master list",
    "build_cache.py            84 lines  - Data caching utility for faster loads",
    "perf_utils.py             66 lines  - Performance optimization helpers",
    "get_4_data.py             45 lines  - Data fetch helper functions",
], 0.35, 4.75, 6.05, 2.37, clr=GREE, tsz=12, isz=10)

feat_card(sl, "Email Notification System (Gmail SMTP)", [
    "Automated HTML email sent on every new claim submission",
    "Email includes: Customer info, product model, OSID, issue details",
    "Supports file attachments (customer invoices / product photos)",
    "Configurable CC list via environment variable",
    "Sender credentials secured via env vars (not hardcoded)",
], 6.55, 1.62, 6.43, 2.5, clr=PINK, tsz=12, isz=10)

feat_card(sl, "Production Deployment  (Render Cloud)", [
    "Procfile        - Gunicorn WSGI server configuration",
    "render.yaml     - Cloud environment + build commands",
    "runtime.txt     - Python version (3.12) pinned for consistency",
    "requirements.txt  - 8 dependencies declared (Flask, psycopg2, gspread...)",
    ".gitignore      - 15 exclusion patterns (secrets, uploads, cache)",
    "PostgreSQL DB   - Provisioned on Render with DATABASE_URL env var",
    "Environment vars: DB_URL, SMTP creds, Sheet ID, WEB_APP_URL, etc.",
], 6.55, 4.28, 6.43, 2.84, clr=CYAN, tsz=12, isz=10)

btm(sl, [
    ("5", "CSS Files"), ("2,725", "CSS Lines"),
    ("274", "Google Script Lines"), ("15.8 MB", "OSID Reference Data"), ("Live", "on Render Cloud"),
], BLUE)


# =============================================================
# SLIDE 4 - MAY 2026
# =============================================================
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
hdr(sl, "May 2026  —  Feature Enhancements & Data Refresh",
    "6 Commits between May 2 and May 26, 2026  |  2 Feature Updates  |  4 Data Files Refreshed", PURP)

feat_card(sl, "May 2, 2026 - Commit b0d1578: Added 'Cancelled' Status End-to-End  (5 files)", [
    "app.py          Added 'Cancelled' to the ALLOWED_STATUSES whitelist (backend validation)",
    "dashboard.html  Added 'Cancelled' as option in Edit Claim status dropdown",
    "analytics.html  Added 'Cancelled' KPI card to the analytics dashboard header",
    "analytics.js    Added Cancelled count query + chart segment + colour code logic",
    "script.js       Added 'Cancelled' to client-side status filter + CSS class mapping",
    "",
    "Result: Staff can now officially mark any claim as Cancelled across all portal sections",
    "Impact: Analytics now shows Cancelled count in real time alongside other statuses",
], 0.35, 1.62, 6.2, 4.72, clr=PURP, tsz=12.5, isz=10.5)

feat_card(sl, "May 4, 2026 - Commit 0d7a0d7: Dashboard Aging Bucket Fix  (2 files, 90 lines)", [
    "app.py          Rewrote aging algorithm (59 lines changed)",
    "                Fixed: 0-7 days / 8-15 days / 16-30 days / 30+ days groupings",
    "                Corrected TAT (Turn Around Time) calculation edge cases",
    "dashboard.html  Removed hardcoded merge cell structure (15 lines fixed)",
    "                Dashboard table now matches the original Excel report design exactly",
    "",
    "Why it mattered: Manager review relies on correct aging buckets for SLA tracking",
], 6.65, 1.62, 6.33, 3.0, clr=AMBE, tsz=12.5, isz=10.5)

feat_card(sl, "May Data File Refreshes  (4 commits)", [
    "OSID Data File   Updated to Feb 2026 data  (Commit: 268ea6f  |  May 17)",
    "EXCEL_FILE ref   Updated to point to Feb 2026 file  (Commit: b43b012  |  May 19)",
    "Future Store List.xlsx   New stores added for Q1 2026  (Commit: 55f1138  |  May 22)",
    "RBM, BDM, Branch.xlsx   Team mapping updated  (Commit: 55f1138  |  May 22)",
    "myG All Store.xlsx       Store master refreshed  (Commit: 55f1138  |  May 22)",
    "app.py  Minor config tweaks  (Commit: 48c0815  |  May 26)",
], 6.65, 4.75, 6.33, 2.37, clr=GREE, tsz=12.5, isz=10.5)

btm(sl, [
    ("6", "Total Commits"), ("1", "New Status Added"),
    ("5", "Files for Status Fix"), ("90", "Lines in Aging Fix"),
    ("4", "Data Files Refreshed"),
], PURP)


# =============================================================
# SLIDE 5 - JUNE 2026
# =============================================================
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
hdr(sl, "June 2026  —  Operations, Monitoring & Live Issue Resolution",
    "Portal fully in production  |  First live claims managed  |  Early issues identified and tracked", GREE)

feat_card(sl, "Production Portal Monitoring", [
    "Monitored live claims being submitted through the portal (real customers)",
    "Tracked claim status progression: Submitted -> Registered -> Follow Up -> Resolved",
    "Verified Google Sheets sync running every 20 seconds (data consistency checks)",
    "Monitored PostgreSQL DB for correct claim data storage and retrieval",
    "Checked portal uptime on Render cloud (production server availability)",
    "Validated email notification delivery on every new claim submission",
    "Tracked Replace Workflow: Customer Confirmation, OSG Approval, Invoice, Settlement",
], 0.35, 1.62, 6.1, 4.72, clr=GREE, tsz=12.5, isz=10.5)

feat_card(sl, "Staff Onboarding & Training", [
    "Onboarded OSG service store staff to the new portal system",
    "Demonstrated claim submission workflow end-to-end",
    "Explained claim status update procedure and follow-up notes",
    "Trained staff on using OSID lookup to auto-fill invoice/serial data",
    "Collected UX feedback: slow workflows, confusing status labels",
], 6.55, 1.62, 6.43, 2.45, clr=AMBE, tsz=12.5, isz=10.5)

feat_card(sl, "Issues Identified in June  (Fixed Later in July)", [
    "Sync conflict: Google Sheets sync overwriting portal-set statuses",
    "Missing Replacement Approved claims from old Google Sheet",
    "Workflow checkboxes not auto-ticking for Replacement Approved claims",
    "Old claims being appended to bottom of new Google Sheet on save",
    "WhatsApp notifications needed cutoff date to avoid spamming old customers",
    "Data identified: Duplicate claim record needed to be split into two",
], 6.55, 4.22, 6.43, 2.9, clr=RED, tsz=12.5, isz=10.5)

btm(sl, [
    ("567+", "Claims in Portal"), ("100%", "Server Uptime"),
    ("20 sec", "Sync Interval"), ("4+", "Issues Identified"), ("Staff", "Onboarded"),
], GREE)


# =============================================================
# SLIDE 6 - KEY CHALLENGES
# =============================================================
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
hdr(sl, "Key Challenges Faced  —  Q2 2026",
    "Technical and operational challenges during April-June 2026 development cycle", PINK)

chal = [
    (PINK, "Building Full-Stack Portal Alone in April (Single Sprint)",
     "Had to design DB schema, backend API, frontend UI, Google Sheets integration, email system and cloud deployment all simultaneously. No existing codebase to start from.",
     "Delivered complete working portal by April 20 — 35 files, 12,776 lines, fully deployed to Render production with PostgreSQL, Gmail SMTP, Google Sheets API all live."),

    (AMBE, "Migrating 500+ Legacy Claims from Google Sheets",
     "Old Google Sheets had inconsistent column names, missing fields, mixed date formats and duplicate entries across hundreds of rows. Could not import directly.",
     "Built ClaimWrapper adapter class mapping old Sheet column headers to new DB schema. Wrote normalisation logic to handle all format variations during import."),

    (PURP, "Dashboard Aging Buckets Not Matching Original Design  (May 4)",
     "The aging bucket report (0-7d / 8-15d / 16-30d / 30+d) was calculating incorrectly. The table layout also didn't match the original Excel merge-cell design managers expected.",
     "Rewrote aging algorithm in app.py (59 lines changed) and fixed dashboard.html structure. Final output now exactly matches original Excel report design specifications."),

    (RED, "Google Sheets Sync Overwriting Portal Statuses  (Identified June)",
     "Every 20-second sync cycle was overwriting manually-set portal statuses (Replacement Approved, Repair Completed etc.) back to earlier values, silently reverting work done by staff.",
     "Identified root cause in June during monitoring. Documented for fix. Resolved in July by adding PROTECTED_STATUSES logic in pg_sync.py to preserve all terminal statuses."),
]

for i, (clr, title, prob, sol) in enumerate(chal):
    col = i % 2; row = i // 2
    x = 0.35 + col * 6.5
    y = 1.62 + row * 2.7
    Rc(sl, x, y, 6.28, 2.58, fill=CARD2, line=BORD, lw=0.3)
    Rc(sl, x, y, 6.28, 0.048, fill=clr)
    Tx(sl, title, x+0.15, y+0.1, 6.0, 0.42, sz=12, bold=True, col=clr)
    Tx(sl, "Problem:", x+0.15, y+0.59, 0.9, 0.27, sz=9.5, bold=True, col=RED)
    bullets(sl, [prob], x+1.05, y+0.59, 5.1, 0.68, sz=9.5, col=MUTED)
    Rc(sl, x+0.15, y+1.38, 5.95, 0.025, fill=BORD)
    Tx(sl, "Resolution:", x+0.15, y+1.49, 1.05, 0.27, sz=9.5, bold=True, col=GREE)
    bullets(sl, [sol], x+1.2, y+1.49, 4.95, 0.92, sz=9.5, col=WHITE)

Rc(sl, 0, 7.42, 13.33, 0.04, fill=PINK)


# =============================================================
# SLIDE 7 - THANK YOU
# =============================================================
sl = prs.slides.add_slide(BLANK)
page_bg(sl)

# Diagonal lines BG
for i in range(22):
    Rc(sl, i * 0.62, 0, 0.025, 7.5, fill=RGBColor(0x0E, 0x1C, 0x35))

for y, c in [(0, BLUE), (0.06, PURP), (0.12, PINK)]:
    Rc(sl, 0, y, 13.33, 0.055, fill=c)

Tx(sl, "Thank You", 0, 1.5, 13.33, 1.55,
   sz=72, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
Tx(sl, "Questions  &  Open Discussion", 0, 3.08, 13.33, 0.68,
   sz=22, col=BLUE, align=PP_ALIGN.CENTER)

Rc(sl, 4.0, 3.95, 5.33, 0.04, fill=PURP)

Tx(sl, "OSG myG Portal   |   Q2 2026 Developer Work Summary   |   Manager Review",
   0, 4.12, 13.33, 0.45, sz=13, italic=True, col=MUTED, align=PP_ALIGN.CENTER)

for i, (mon, val, lbl, clr) in enumerate([
    ("APRIL",  "12,776",   "Lines of Code Written",   BLUE),
    ("MAY",    "6",        "Commits / Enhancements",   PURP),
    ("JUNE",   "567+",     "Claims in Production",     GREE),
]):
    x = 1.8 + i * 3.3
    Rc(sl, x, 5.1, 3.0, 1.75, fill=CARD, line=BORD, lw=0.3)
    Rc(sl, x, 5.1, 3.0, 0.05, fill=clr)
    Tx(sl, mon, x, 5.22, 3.0, 0.4, sz=13, bold=True, col=clr, align=PP_ALIGN.CENTER)
    Tx(sl, val, x, 5.62, 3.0, 0.75, sz=30, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    Tx(sl, lbl, x, 6.35, 3.0, 0.38, sz=10, col=MUTED, align=PP_ALIGN.CENTER)

for y, c in [(7.4, PINK), (7.45, PURP), (7.5, BLUE)]:
    Rc(sl, 0, y, 13.33, 0.05, fill=c)


out = "OSG_Portal_Q2_2026_FINAL.pptx"
prs.save(out)
print(f"[DONE] Saved: {out}")
