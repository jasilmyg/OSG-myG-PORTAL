from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Function to add a title slide
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

# Function to add a bullet point slide
def add_bullet_slide(prs, title, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = title
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        if isinstance(point, list):
            for subpoint in point:
                p = tf.add_paragraph()
                p.text = subpoint
                p.level = 1

# Slide 1: Title
add_title_slide(
    prs, 
    "Q2 2026 Accomplishments", 
    "Review of Development Work\nApril, May & June 2026"
)

# Slide 2: Executive Summary
add_bullet_slide(
    prs,
    "Executive Summary",
    [
        "Successfully developed and deployed the Enterprise Retail Analytics Dashboard.",
        "Integrated Advanced AI and Machine Learning forecasting algorithms (Scikit-Learn, Deep Learning).",
        "Optimized massive database aggregations processing 12.6+ million rows for near-instant rendering.",
        "Delivered the automated 'SHE START' live-syncing dashboard with Google Sheets.",
        "Significantly improved the stability, speed, and intelligence of the main loyalty portal."
    ]
)

# Slide 3: Enterprise Retail Analytics Dashboard
add_bullet_slide(
    prs,
    "Enterprise Retail Analytics Dashboard",
    [
        "Overhauled the UI into an Enterprise-grade Business Intelligence Portal without the need for manual uploads.",
        "Implemented high-speed DataTables grids for month-wise, quarterly, and FY comparisons.",
        "Added Dynamic Report Generators allowing users to export data to custom-formatted Excel (using Calamine engine for 10x speed) and PDF.",
        "Engineered Average Selling Price (ASP) and category scorecards that map directly to business logic.",
        "Resolved complex date parsing bugs that previously corrupted data (American vs. European date formats)."
    ]
)

# Slide 4: AI & Machine Learning Integrations
add_bullet_slide(
    prs,
    "Advanced AI & Machine Learning Integrations",
    [
        "Developed a Futuristic Neural Intelligence System for Customer Campaign Analysis.",
        "Implemented Scikit-Learn models (Random Forest, MLPRegressor, GradientBoostingRegressor) for live predictive scoring and forecasting.",
        "Built the 'MalayalamCalendarFeaturizer' to generate advanced seasonal vectors (e.g., Days to Onam) for precise localized predictions.",
        "Created an Advanced AI Insights Engine that dynamically analyzes data arrays and generates JSON-driven insights, deep analysis, and recommendations.",
        "Designed UI elements like Probability Gauges, Semi-Donuts, and Dormancy Risk Meters."
    ]
)

# Slide 5: Database & Performance Optimizations
add_bullet_slide(
    prs,
    "Database & Performance Optimizations",
    [
        "Refactored sequential database queries into high-performance Django ORM aggregations.",
        "Integrated Django's high-speed Memory Caching system (Redis) to deliver massive datasets in milliseconds.",
        "Fixed cross-join duplications in materialized views (e.g., 'mv_yearly_cohort') to ensure 100% accuracy in Cohort Retention and LTV calculation.",
        "Built pre-aggregation engines ('mv_dormant_reactivation') to accurately bucket millions of historical records.",
        "Optimized Excel parsing from openpyxl to the Rust-based 'calamine' engine to handle massive 100MB+ data files efficiently."
    ]
)

# Slide 6: SHE START Applicant Tracking Dashboard
add_bullet_slide(
    prs,
    "SHE START Applicant Tracking System",
    [
        "Developed a live Google Sheets sync engine (via gspread) with a 25-second auto-refresh interval.",
        "Implemented interactive inline editing allowing panelists to input scores that silently save to the local Postgres database.",
        "Created a complex Final Decision Criteria matrix, mathematically dropping the highest and lowest scores across 6 panelists to find the true weighted average.",
        "Automated decision badges (Strong Final Selection, Recommended for Top 10, Waitlist).",
        "Configured secure, role-based access control to completely isolate the SHE START views from normal users."
    ]
)

# Slide 7: Technical Stack & Architecture
add_bullet_slide(
    prs,
    "Technical Architecture & Tools",
    [
        "Backend: Python, Django, PostgreSQL, Pandas, Numpy, Calamine.",
        "AI/ML: Scikit-Learn (Random Forest, MLPRegressor), PyTorch concepts.",
        "Frontend: HTML5, JavaScript, Plotly (Charts), DataTables, Custom CSS.",
        "Infrastructure: Render Deployment, SSH Tunnels, Git Version Control, Environment Variables.",
        "Integrations: Google Sheets API (gspread), Advanced Excel Generation."
    ]
)

# Save the presentation
output_path = "c:\\Users\\jasil_myg\\Desktop\\OSG-myG-PORTAL-mainnnnn - Copy\\April_May_June_2026_Work_Summary.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
