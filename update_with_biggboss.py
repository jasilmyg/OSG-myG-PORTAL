from pptx import Presentation

# 1. Update the Markdown Report
md_file = r"C:\Users\jasil_myg\.gemini\antigravity-ide\brain\2d556046-4336-4550-afe9-c59e03f7448f\Detailed_Work_Report_Q2_2026.md"
with open(md_file, 'a', encoding='utf-8') as f:
    f.write("\n## 8. BIGG BOSS & FoneFlix Registration Portals\n")
    f.write("Built a highly customized Flask-based web application initially designed for 'Bigg Boss Season 8 – Agnipareeksha' auditions, which was later successfully pivoted into the 'myG FoneFlix Mobile Phone Short Film Contest 2026'.\n\n")
    f.write("* **Video Upload Architecture (Google Drive OAuth):**\n")
    f.write("  * Engineered a custom OAuth 2.0 Client ID integration that allowed users to authenticate and upload large video files (auditions / short films) directly into their personal Google Drive folders, bypassing the Render server's ephemeral storage limits.\n")
    f.write("* **Dynamic UI/UX Design:**\n")
    f.write("  * Designed a complex, responsive hero section using precise grid layouts (40/60 split) matching OTT reality-show aesthetics.\n")
    f.write("  * Implemented neon-glow typography, dark glassmorphism form backgrounds, and floating 3D particle animations for brand elements.\n")
    f.write("* **Project Pivot & Rebranding:**\n")
    f.write("  * Successfully transitioned the entire codebase from the Bigg Boss branding (Navy Blue & Neon Pink) to the FoneFlix cinematic branding (Dark Browns & Orange).\n")
    f.write("  * Refactored form inputs, consent clauses, and loading UI states to match the new short film contest requirements.\n")


# 2. Update the PPTX
pptx_path = r"c:\Users\jasil_myg\Desktop\OSG-myG-PORTAL-mainnnnn - Copy\OSG_Portal_WorkSummary_Q2_2026_Detailed.pptx"
prs = Presentation(pptx_path)

bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = "Bigg Boss & FoneFlix Registration Portals"

body_shape = shapes.placeholders[1]
tf = body_shape.text_frame

points = [
    "Developed a Flask-based registration portal for 'Bigg Boss Season 8 – Agnipareeksha' and pivoted it into the 'myG FoneFlix Short Film Contest'.",
    "Engineered custom Google Drive OAuth 2.0 integration to allow large video uploads directly to user's Drive (bypassing server storage limits).",
    "Designed an ultra-premium OTT streaming UI with glassmorphism, neon glows, and floating 3D particle animations.",
    "Completely rebranded the platform dynamically from a neon-pink Bigg Boss theme to a cinematic orange/brown FoneFlix theme."
]

tf.text = points[0]
for point in points[1:]:
    p = tf.add_paragraph()
    p.text = point
    p.level = 0

prs.save(pptx_path)
print("Updated Markdown and PPTX successfully.")
